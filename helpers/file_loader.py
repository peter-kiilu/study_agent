import os
from pypdf import PdfReader
from pptx import Presentation
from docx import Document


# ── Supported extensions ─────────────────────────────────────────────────────
SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".pptx", ".docx"}


def load_notes(file_path: str) -> str:
    """
    Load lecture notes from a supported file.
    Supported: .txt, .pdf, .pptx, .docx
    Returns the raw text content.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _load_pdf(file_path)
    elif ext == ".txt":
        return _load_txt(file_path)
    elif ext == ".pptx":
        return _load_pptx(file_path)
    elif ext == ".docx":
        return _load_docx(file_path)
    else:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Use {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )


# ── Individual loaders ───────────────────────────────────────────────────────

def _load_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text.strip()


def _load_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read().strip()


def _load_pptx(file_path: str) -> str:
    """Extract text from all slides in a PowerPoint file."""
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        text_parts.append(line)
    return "\n".join(text_parts)


def _load_docx(file_path: str) -> str:
    """Extract text from all paragraphs in a Word document."""
    doc = Document(file_path)
    text_parts = []
    for para in doc.paragraphs:
        line = para.text.strip()
        if line:
            text_parts.append(line)
    return "\n".join(text_parts)


# ── Directory scanner ────────────────────────────────────────────────────────

def scan_directory(directory_path: str) -> dict:
    """
    Recursively scan a directory for supported note files.

    Returns a dict mapping course folder names to lists of file info:
        {
            "Course Name": [
                {"path": "/full/path/to/file.pdf", "name": "file.pdf", "ext": ".pdf"},
                ...
            ],
            ...
        }

    Top-level files (not inside a subfolder) are grouped under "General".
    """
    if not os.path.isdir(directory_path):
        raise NotADirectoryError(f"Not a directory: {directory_path}")

    courses = {}

    for entry in sorted(os.listdir(directory_path)):
        full_path = os.path.join(directory_path, entry)

        if os.path.isdir(full_path):
            # Each subfolder is treated as a course
            files = _collect_files(full_path)
            if files:
                courses[entry] = files

        elif os.path.isfile(full_path):
            ext = os.path.splitext(entry)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                courses.setdefault("General", []).append({
                    "path": full_path,
                    "name": entry,
                    "ext": ext,
                })

    return courses


def _collect_files(folder_path: str) -> list:
    """Collect all supported files recursively within a folder."""
    files = []
    for root, _dirs, filenames in os.walk(folder_path):
        for fname in sorted(filenames):
            ext = os.path.splitext(fname)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                files.append({
                    "path": os.path.join(root, fname),
                    "name": fname,
                    "ext": ext,
                })
    return files


def load_all_from_directory(directory_path: str) -> tuple:
    """
    Scan a directory, load ALL discovered note files, and return
    (courses_dict, combined_text).

    courses_dict is the output of scan_directory() (for display).
    combined_text is every file's contents joined together, with headers.
    """
    courses = scan_directory(directory_path)

    if not courses:
        raise FileNotFoundError(
            f"No supported note files found in: {directory_path}\n"
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    text_parts = []
    for course_name, files in courses.items():
        text_parts.append(f"\n{'=' * 60}")
        text_parts.append(f"COURSE: {course_name}")
        text_parts.append(f"{'=' * 60}\n")

        for file_info in files:
            try:
                content = load_notes(file_info["path"])
                if content:
                    text_parts.append(f"--- {file_info['name']} ---")
                    text_parts.append(content)
                    text_parts.append("")  # blank line separator
            except Exception as e:
                # Skip files that fail to load, but note them
                text_parts.append(f"--- {file_info['name']} [SKIPPED: {e}] ---")

    return courses, "\n".join(text_parts).strip()
